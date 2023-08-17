import os
from pptx import Presentation
from transformers import BartTokenizer, BartForConditionalGeneration

def extract_slides_from_pptx(file_path):
    presentation = Presentation(file_path)
    slides_text = []

    for slide in presentation.slides:
        slide_text = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)

        slides_text.append('\n'.join(slide_text))

    return slides_text

def generate_summary(text):
    model_name = "facebook/bart-large-cnn"
    tokenizer = BartTokenizer.from_pretrained(model_name)
    model = BartForConditionalGeneration.from_pretrained(model_name)

    inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs, max_length=75, min_length=30, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)

    return summary

# Root directory containing the pptx files
root_dir = 'input/content'

# Find pptx files in a root_dir and read them
def find_and_read_ppt_files(root_dir):
    pptx_files = []

    for root, dirs, files in os.walk(root_dir):
        for file in files:
            if file.endswith(".pptx"):
                pptx_path = os.path.join(root, file)
                pptx_files.append(pptx_path)

    return pptx_files

pptx_files = find_and_read_ppt_files(root_dir)

for pptx_path in pptx_files:
    slides_text = extract_slides_from_pptx(pptx_path)

    # Extract the base filename without extension
    base_filename = os.path.splitext(os.path.basename(pptx_path))[0]

    # Create a directory for slide summaries
    summary_dir = os.path.join(os.path.dirname(pptx_path), f"{base_filename}_summary")
    os.makedirs(summary_dir, exist_ok=True)

    for i, slide_text in enumerate(slides_text):
        summary = generate_summary(slide_text)

        # Save summary in the slide summary subfolder
        summary_filename = f"{base_filename}_summary_slide_{i + 1}.comment"
        summary_path = os.path.join(summary_dir, summary_filename)

        # Overwrite existing summary file if it exists
        if os.path.exists(summary_path):
            os.remove(summary_path)

        with open(summary_path, 'w', encoding='utf-8') as summary_file:
            summary_file.write(summary)

        print(f"Summary saved for slide {i + 1} in {pptx_path} at {summary_path}")

